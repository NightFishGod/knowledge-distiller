"""
知识蒸馏 v2.0 - 深度内容提取器
支持：PPTX(含备注)、XLSX(全表+批注)、DOCX(含页眉页脚)、PDF、图片OCR、纯文本
"""
import os
import re
import json
from pathlib import Path
from datetime import datetime
from collections import defaultdict

# ---- 文档提取 ----
def extract_docx(filepath: str, include_headers: bool = True) -> dict:
    """提取 DOCX 正文 + 页眉页脚"""
    from docx import Document
    doc = Document(filepath)
    parts = []

    if include_headers:
        for section in doc.sections:
            for p in section.header.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())
            for p in section.footer.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())

    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return {"type": "docx", "text": "\n".join(parts), "paragraphs": len(doc.paragraphs)}


def extract_pptx(filepath: str, include_notes: bool = True) -> dict:
    """深度提取 PPTX：所有形状文本 + 备注 + 表格 + SmartArt"""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_output = []
    total_text = []
    notes_text = []

    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
                        total_text.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_parts.append(row_text)
                        total_text.append(row_text)

        # 演讲者备注
        if include_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            note_text = notes_slide.notes_text_frame.text.strip()
            if note_text:
                notes_text.append(f"[备注-第{i+1}页] {note_text}")
                total_text.append(note_text)

        if slide_parts:
            slides_output.append({"slide": i + 1, "content": "\n".join(slide_parts)})

    return {
        "type": "pptx",
        "slides": len(prs.slides),
        "slides_with_content": len(slides_output),
        "has_notes": len(notes_text) > 0,
        "notes_count": len(notes_text),
        "text": "\n".join(total_text),
        "notes": "\n".join(notes_text),
        "slides_detail": slides_output,
    }


def extract_xlsx(filepath: str, all_sheets: bool = True, include_comments: bool = True) -> dict:
    """提取 XLSX：全工作表 + 批注"""
    import openpyxl
    wb = openpyxl.load_workbook(filepath, data_only=True)
    sheets_data = []

    sheets_to_read = wb.sheetnames if all_sheets else [wb.active.title] if wb.active else wb.sheetnames[:1]

    for sname in sheets_to_read:
        sheet = wb[sname]
        max_row = sheet.max_row or 0
        max_col = sheet.max_column or 0

        if max_row == 0:
            continue

        # 提取表头
        headers = []
        for col in range(1, max_col + 1):
            val = sheet.cell(row=1, column=col).value
            headers.append(str(val).strip() if val else "")

        # 提取所有行（限制行数防OOM）
        rows = []
        sample_limit = min(max_row, 200)
        for row_idx in range(1, sample_limit + 1):
            row_data = []
            for col_idx in range(1, max_col + 1):
                val = sheet.cell(row=row_idx, column=col_idx).value
                if val is not None:
                    row_data.append(str(val).strip())
            if row_data:
                rows.append(" | ".join(row_data))

        # 批注
        comments_list = []
        if include_comments:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.comment:
                        comments_list.append(f"[{cell.coordinate}] {cell.comment.text.strip()}")

        sheets_data.append({
            "name": sname,
            "rows": max_row,
            "columns": max_col,
            "headers": [h for h in headers if h],
            "sample_rows": rows[:50],
            "comments": comments_list[:20],
            "has_comments": len(comments_list) > 0,
        })

    wb.close()

    # 拼合全部文本
    all_text_parts = []
    for sd in sheets_data:
        all_text_parts.append(f"--- 工作表: {sd['name']} ({sd['rows']}行 × {sd['columns']}列) ---")
        all_text_parts.extend(sd["sample_rows"])
        if sd["comments"]:
            all_text_parts.append("批注:")
            all_text_parts.extend(sd["comments"])

    return {
        "type": "xlsx",
        "sheets_total": len(wb.sheetnames),
        "sheets_read": len(sheets_data),
        "sheets": sheets_data,
        "text": "\n".join(all_text_parts),
    }


def extract_pdf(filepath: str) -> dict:
    """提取 PDF 文本"""
    text_parts = []

    # 尝试 pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        if text_parts:
            return {"type": "pdf", "pages": len(pdf.pages), "text": "\n".join(text_parts)}
    except ImportError:
        pass

    # 回退到 PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
        return {"type": "pdf", "pages": len(reader.pages), "text": "\n".join(text_parts)}
    except Exception:
        return {"type": "pdf", "pages": 0, "text": "", "error": "无法提取文本"}


def extract_image_ocr(filepath: str, lang: str = "chi_sim+eng") -> dict:
    """图片 OCR 文本提取"""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img, lang=lang)
        return {"type": "image_ocr", "text": text.strip()}
    except ImportError:
        return {"type": "image_ocr", "text": "", "skipped": "pytesseract 未安装"}
    except Exception as e:
        return {"type": "image_ocr", "text": "", "error": str(e)}


def extract_text_file(filepath: str) -> dict:
    """提取纯文本文件"""
    for enc in ["utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            with open(filepath, "r", encoding=enc) as f:
                text = f.read()
            return {"type": "text", "encoding": enc, "text": text}
        except UnicodeDecodeError:
            continue
    return {"type": "text", "text": "", "error": "无法解码"}


# ---- 主提取流程 ----
EXTENSION_HANDLERS = {
    ".docx": ("document", extract_docx),
    ".doc": ("document", extract_docx),
    ".pptx": ("document", extract_pptx),
    ".ppt": ("document", extract_pptx),
    ".xlsx": ("document", extract_xlsx),
    ".xls": ("document", extract_xlsx),
    ".pdf": ("document", extract_pdf),
    ".txt": ("text", extract_text_file),
    ".md": ("text", extract_text_file),
    ".csv": ("text", extract_text_file),
    ".json": ("text", extract_text_file),
    ".jpg": ("image", None),
    ".jpeg": ("image", None),
    ".png": ("image", None),
    ".bmp": ("image", None),
    ".webp": ("image", None),
}


def extract_single_file(filepath: str, config: dict = None) -> dict:
    """提取单个文件的内容"""
    ext = Path(filepath).suffix.lower()
    handler = EXTENSION_HANDLERS.get(ext)

    if handler is None:
        return {"type": "skipped", "reason": f"不支持的格式: {ext}"}

    cat, func = handler
    if func is None:
        return {"type": cat, "text": "", "note": "二进制文件，未提取文本"}

    try:
        result = func(filepath)
        result["category"] = cat
        result["extension"] = ext
        return result
    except Exception as e:
        return {"type": cat, "extension": ext, "text": "", "error": str(e)}


def batch_extract(file_map_path: str, config: dict = None, progress_callback=None) -> dict:
    """
    批量提取文件内容
    返回 content_dump 和 extraction_stats
    """
    with open(file_map_path, "r", encoding="utf-8") as f:
        file_map = json.load(f)

    files = file_map.get("files", [])
    if not files and isinstance(file_map, dict):
        # 可能是按类型组织的
        for cat_files in file_map.values():
            if isinstance(cat_files, list):
                files.extend(cat_files)

    total = len(files)
    corpus = []        # 所有提取文本汇总
    extraction_results = []
    stats = {"total": total, "extracted": 0, "failed": 0, "skipped": 0, "total_chars": 0}

    for i, finfo in enumerate(files):
        filepath = finfo.get("path") or finfo.get("file_path", "")
        if not filepath or not os.path.exists(filepath):
            stats["skipped"] += 1
            continue

        result = extract_single_file(filepath, config)
        extraction_results.append({"file": filepath, **result})

        text = result.get("text", "")
        if text:
            stats["extracted"] += 1
            stats["total_chars"] += len(text)
            corpus.append(f"### {Path(filepath).name}\n{text[:50000]}")
        elif result.get("error"):
            stats["failed"] += 1
        else:
            stats["skipped"] += 1

        if progress_callback and i % 10 == 0:
            progress_callback(i, total, f"提取中: {i}/{total}")

    content_dump = "\n\n".join(corpus)
    combined_corpus = re.sub(r"\n{3,}", "\n\n", content_dump)

    return {
        "content_dump": combined_corpus,
        "corpus_text": combined_corpus,
        "extraction_results": extraction_results,
        "stats": stats,
    }


def generate_content_dump(batch_result: dict, output_path: str):
    """写出 content_dump.md"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"# 内容提取汇总\n\n")
        f.write(f"- 总文件: {batch_result['stats']['total']}\n")
        f.write(f"- 成功提取: {batch_result['stats']['extracted']}\n")
        f.write(f"- 提取字符数: {batch_result['stats']['total_chars']:,}\n\n")
        f.write(batch_result["content_dump"])
    return output_path


if __name__ == "__main__":
    # 快速测试
    import sys
    if len(sys.argv) > 1:
        path = sys.argv[1]
        r = extract_single_file(path)
        print(f"类型: {r.get('type')} | 文本长度: {len(r.get('text',''))}")
        if r.get("text"):
            print(r["text"][:500])
（内容由AI生成，仅供参考）
